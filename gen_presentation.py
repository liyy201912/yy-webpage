import sys

sys.path.insert(0, "/Users/yli16/projects/.pip_pkgs")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
ACCENT_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
ACCENT_GREEN = RGBColor(0x2F, 0x85, 0x59)
ACCENT_PURPLE = RGBColor(0x73, 0x4A, 0xB8)
ACCENT_ORANGE = RGBColor(0xD9, 0x62, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xF4, 0xF7, 0xFB)
CARD_BG = RGBColor(0xFA, 0xFB, 0xFD)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_GRAY = RGBColor(0x6B, 0x72, 0x80)
NEAR_BLACK = RGBColor(0x1F, 0x29, 0x37)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_font(run, size=18, bold=False, color=NEAR_BLACK, name="Helvetica Neue"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_paragraph(tf, text, size=18, bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT, space_after=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return p


def add_title(slide, title, subtitle=None, accent=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.08), Inches(0.48))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = add_textbox(slide, Inches(1.08), Inches(0.42), Inches(11.2), Inches(0.65))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=30, bold=True, color=DARK_BLUE)

    if subtitle:
        add_paragraph(tf, subtitle, size=15, color=LIGHT_GRAY, space_after=Pt(0))


def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_metric(slide, left, top, width, height, number, label, accent=ACCENT_BLUE):
    shape = add_card(slide, left, top, width, height, fill_color=SOFT_BG, border_color=SOFT_BG)
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)
    r = tf.paragraphs[0].add_run()
    r.text = number
    set_font(r, size=26, bold=True, color=accent)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    set_font(r2, size=12, color=LIGHT_GRAY)


def add_bullet_panel(slide, left, top, width, height, title, bullets, accent=ACCENT_BLUE):
    add_card(slide, left, top, width, height)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = add_textbox(slide, left + Inches(0.24), top + Inches(0.18), width - Inches(0.35), height - Inches(0.28))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=19, bold=True, color=DARK_BLUE)
    p.space_after = Pt(8)

    for bullet in bullets:
        add_paragraph(tf, "• " + bullet, size=13, color=LIGHT_GRAY, space_after=Pt(5))


def add_pub_card(slide, left, top, width, height, title, venue, cites, desc, accent=ACCENT_BLUE):
    add_card(slide, left, top, width, height, fill_color=SOFT_BG, border_color=SOFT_BG)
    tb = add_textbox(slide, left + Inches(0.22), top + Inches(0.12), width - Inches(0.34), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=15, bold=True, color=DARK_BLUE)
    p.space_after = Pt(2)
    add_paragraph(tf, desc, size=12, color=LIGHT_GRAY, space_after=Pt(6))
    add_paragraph(tf, f"{venue}  |  {cites}", size=11, bold=True, color=accent, space_after=Pt(0))


def add_simple_list(slide, left, top, width, height, title, items, accent=ACCENT_BLUE):
    add_card(slide, left, top, width, height)
    tb = add_textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.34), height - Inches(0.28))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, size=18, bold=True, color=DARK_BLUE)
    p.space_after = Pt(8)
    for item in items:
        add_paragraph(tf, item, size=12, color=LIGHT_GRAY, space_after=Pt(5))


# Slide 1: title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BLUE)

tb = add_textbox(slide, Inches(1.3), Inches(1.75), Inches(10.8), Inches(1.25))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Yanyu Li"
set_font(r, size=50, bold=True, color=WHITE)
add_paragraph(tf, "Ph.D., Northeastern University", size=23, color=RGBColor(0xC8, 0xD8, 0xEE), align=PP_ALIGN.CENTER, space_after=Pt(18))
add_paragraph(
    tf,
    "Efficient AI  |  Vision Transformers  |  Diffusion Models  |  Mobile and Edge Deployment",
    size=17,
    color=RGBColor(0xA6, 0xBB, 0xD6),
    align=PP_ALIGN.CENTER,
    space_after=Pt(14),
)
add_paragraph(tf, "59 publications  |  2,963 citations  |  h-index 23", size=15, color=RGBColor(0x92, 0xAB, 0xC9), align=PP_ALIGN.CENTER, space_after=Pt(0))

tb2 = add_textbox(slide, Inches(2.2), Inches(5.65), Inches(9.0), Inches(0.6))
tf2 = tb2.text_frame
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "yanyu.li@northeastern.edu   |   scholar.google.com/citations?user=XUj8koUAAAAJ   |   github.com/liyy201912"
set_font(r, size=13, color=RGBColor(0xA6, 0xBB, 0xD6))


# Slide 2: profile overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Research Profile", "Scientist working on efficient deep learning for vision, generative modeling, and deployment.")

tb = add_textbox(slide, Inches(1.0), Inches(1.35), Inches(7.1), Inches(1.65))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = (
    "My research focuses on translating state-of-the-art machine learning systems to mobile, edge, "
    "and hardware-constrained environments through efficient architecture design, compression, and "
    "algorithm-hardware co-design."
)
set_font(r, size=18, color=NEAR_BLACK)
p.space_after = Pt(12)
add_paragraph(tf, "Research areas: efficient vision transformers, diffusion models, pruning, quantization, NAS, and on-device generative AI.", size=15, color=LIGHT_GRAY, space_after=Pt(10))
add_paragraph(tf, "Publications at NeurIPS, CVPR, ICCV, ICLR, AAAI, IJCAI, ECCV, HPCA, FPGA, and FPL.", size=15, color=LIGHT_GRAY, space_after=Pt(0))

metrics = [("2,963", "Citations"), ("23", "h-index"), ("59", "Publications"), ("7", "Patents")]
for idx, (num, label) in enumerate(metrics):
    x = Inches(8.65) + (idx % 2) * Inches(2.15)
    y = Inches(1.45) + (idx // 2) * Inches(1.25)
    add_metric(slide, x, y, Inches(1.95), Inches(1.05), num, label)

add_simple_list(
    slide,
    Inches(1.0),
    Inches(3.35),
    Inches(5.3),
    Inches(2.55),
    "Education",
    [
        "Ph.D. in Computer Engineering, Northeastern University (2019-2024)",
        "Advisor: Prof. Yanzhi Wang",
        "M.S. in Mechanical Engineering (Robotics), Boston University (2017-2019)",
        "B.E. in Mechanical Engineering, Tsinghua University (2011-2015)",
    ],
)
add_simple_list(
    slide,
    Inches(6.55),
    Inches(3.35),
    Inches(3.15),
    Inches(2.55),
    "Experience",
    [
        "Snap Inc. research internships (2022-present)",
        "CoCoPIE LLC (2021)",
        "Kuaishou Technology US R&D (2020-2021)",
    ],
    accent=ACCENT_GREEN,
)
add_simple_list(
    slide,
    Inches(9.95),
    Inches(3.35),
    Inches(2.35),
    Inches(2.55),
    "Focus",
    [
        "Mobile AI",
        "Generative models",
        "Efficient training",
        "Edge systems",
    ],
    accent=ACCENT_PURPLE,
)


# Slide 3: research program
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Research Program", "Three connected pillars that link efficient model design to practical deployment.")

add_bullet_panel(
    slide,
    Inches(1.0),
    Inches(1.45),
    Inches(3.75),
    Inches(4.95),
    "Efficient Vision Transformers",
    [
        "Transformer backbones tailored for mobile-scale latency and memory budgets.",
        "Design principles validated through EfficientFormer and EfficientFormerV2.",
        "Applications extend to classification, segmentation, and edge perception.",
    ],
    accent=ACCENT_BLUE,
)
add_bullet_panel(
    slide,
    Inches(4.95),
    Inches(1.45),
    Inches(3.75),
    Inches(4.95),
    "Diffusion Models and Generative AI",
    [
        "Fast text-to-image and video generation for mobile and edge devices.",
        "Quality-control techniques via text encoders, distillation, and efficient sampling.",
        "Compression-aware generative modeling including ultra-low-bit diffusion.",
    ],
    accent=ACCENT_GREEN,
)
add_bullet_panel(
    slide,
    Inches(8.9),
    Inches(1.45),
    Inches(3.45),
    Inches(4.95),
    "Compression and Hardware-Aware Optimization",
    [
        "Quantization, pruning, sparse training, and neural architecture search.",
        "Algorithm-hardware co-design for phones, edge accelerators, and FPGAs.",
        "Real-world deployment emphasis: latency, memory, and energy efficiency.",
    ],
    accent=ACCENT_PURPLE,
)


# Slide 4: efficient vision transformers
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Selected Contributions: Efficient Vision Transformers", "Representative first-author work on mobile-speed visual transformers and efficient deployment.")

add_bullet_panel(
    slide,
    Inches(1.0),
    Inches(1.45),
    Inches(3.1),
    Inches(5.0),
    "Key Outcomes",
    [
        "Showed that transformer backbones can achieve MobileNet-class latency on mobile hardware.",
        "Enabled approximately 1 ms image recognition on smartphones in industrial deployment settings.",
        "Extended efficient design ideas to segmentation and search-based architecture optimization.",
    ],
    accent=ACCENT_BLUE,
)
add_pub_card(
    slide,
    Inches(4.35),
    Inches(1.55),
    Inches(3.7),
    Inches(2.0),
    "EfficientFormer: Vision Transformers at MobileNet Speed",
    "NeurIPS 2022",
    "764 citations",
    "Dimension-consistent transformer design for strong accuracy with mobile-scale inference speed.",
)
add_pub_card(
    slide,
    Inches(8.3),
    Inches(1.55),
    Inches(3.7),
    Inches(2.0),
    "Rethinking Vision Transformers for MobileNet Size and Speed",
    "ICCV 2023",
    "428 citations",
    "Refined mobile transformer design through better accuracy-latency trade-offs and scaling behavior.",
)
add_pub_card(
    slide,
    Inches(4.35),
    Inches(3.85),
    Inches(3.7),
    Inches(1.7),
    "Pruning-as-Search",
    "IJCAI 2022",
    "78 citations",
    "Efficient neural architecture search via channel pruning and structural reparameterization.",
)
add_pub_card(
    slide,
    Inches(8.3),
    Inches(3.85),
    Inches(3.7),
    Inches(1.7),
    "Towards Real-Time Segmentation on the Edge",
    "AAAI 2023",
    "25 citations",
    "Real-time segmentation methods designed for edge deployment scenarios.",
    accent=ACCENT_BLUE,
)


# Slide 5: diffusion and generative AI
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Selected Contributions: Diffusion Models and Generative AI", "Efficient text-to-image and video generation with an emphasis on mobile deployment and quality.")

add_pub_card(
    slide,
    Inches(1.0),
    Inches(1.45),
    Inches(3.85),
    Inches(1.7),
    "SnapFusion",
    "NeurIPS 2023",
    "305 citations",
    "Sub-two-second text-to-image diffusion on mobile devices through architecture and sampling optimization.",
    accent=ACCENT_GREEN,
)
add_pub_card(
    slide,
    Inches(5.0),
    Inches(1.45),
    Inches(3.85),
    Inches(1.7),
    "TextCraftor",
    "CVPR 2024",
    "32 citations",
    "Used the text encoder as a controllable quality signal for diffusion image generation.",
    accent=ACCENT_GREEN,
)
add_pub_card(
    slide,
    Inches(9.0),
    Inches(1.45),
    Inches(3.35),
    Inches(1.7),
    "BitsFusion",
    "NeurIPS 2024",
    "43 citations",
    "1.99-bit diffusion model quantization with limited generation quality degradation.",
    accent=ACCENT_GREEN,
)
add_pub_card(
    slide,
    Inches(1.0),
    Inches(3.45),
    Inches(3.85),
    Inches(1.7),
    "SF-V",
    "NeurIPS 2024",
    "31 citations",
    "Single-forward video generation to reduce computational cost while maintaining quality.",
    accent=ACCENT_GREEN,
)
add_pub_card(
    slide,
    Inches(5.0),
    Inches(3.45),
    Inches(3.85),
    Inches(1.7),
    "LazyDiT",
    "AAAI 2025",
    "49 citations",
    "Accelerated diffusion transformers through lazy learning and reduced redundant computation.",
    accent=ACCENT_GREEN,
)
add_bullet_panel(
    slide,
    Inches(9.0),
    Inches(3.35),
    Inches(3.35),
    Inches(2.1),
    "Broader Generative AI Theme",
    [
        "High-resolution mobile generation (SnapGen, CVPR 2025).",
        "Mobile video generation (SnapGen-V, CVPR 2025).",
        "Recent work also spans efficient autoencoders and diffusion acceleration.",
    ],
    accent=ACCENT_GREEN,
)


# Slide 6: hardware-aware optimization
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Selected Contributions: Compression and Hardware-Aware Optimization", "Quantization, sparse training, and FPGA-aware acceleration for efficient AI systems.")

add_pub_card(
    slide,
    Inches(1.0),
    Inches(1.45),
    Inches(3.85),
    Inches(1.8),
    "Mix and Match",
    "HPCA 2021",
    "156 citations",
    "FPGA-centric mixed-scheme quantization framework for efficient DNN acceleration.",
    accent=ACCENT_PURPLE,
)
add_pub_card(
    slide,
    Inches(5.0),
    Inches(1.45),
    Inches(3.85),
    Inches(1.8),
    "FILM-QNN",
    "FPGA 2022",
    "128 citations",
    "Intra-layer mixed-precision quantization for efficient FPGA deployment of deep networks.",
    accent=ACCENT_PURPLE,
)
add_pub_card(
    slide,
    Inches(9.0),
    Inches(1.45),
    Inches(3.35),
    Inches(1.8),
    "Auto-ViT-Acc",
    "FPL 2022",
    "104 citations",
    "FPGA-aware acceleration framework for vision transformers with mixed-scheme quantization.",
    accent=ACCENT_PURPLE,
)
add_pub_card(
    slide,
    Inches(1.0),
    Inches(3.65),
    Inches(3.85),
    Inches(1.8),
    "Layer Freezing and Data Sieving",
    "NeurIPS 2022",
    "36 citations",
    "Sparse training framework that combines layer freezing with data sieving for efficiency.",
    accent=ACCENT_PURPLE,
)
add_bullet_panel(
    slide,
    Inches(5.0),
    Inches(3.55),
    Inches(7.35),
    Inches(2.0),
    "Systems Perspective",
    [
        "Compression methods are developed with deployment constraints in mind: latency, memory, and energy.",
        "The same optimization perspective carries into diffusion model acceleration and on-device video generation.",
        "Work spans algorithm design, mixed precision, architecture search, and reconfigurable hardware acceleration.",
    ],
    accent=ACCENT_PURPLE,
)


# Slide 7: impact
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Research Impact and Representative Publications", "Selected paper influence from the public Google Scholar profile (March 2026).")

top_pubs = [
    "EfficientFormer (NeurIPS 2022) — 764 citations",
    "EfficientFormerV2 (ICCV 2023) — 428 citations",
    "SnapFusion (NeurIPS 2023) — 305 citations",
    "Mix and Match (HPCA 2021) — 156 citations",
    "FILM-QNN (FPGA 2022) — 128 citations",
    "Auto-ViT-Acc (FPL 2022) — 104 citations",
    "HyperHuman (ICLR 2024) — 83 citations",
    "Pruning-as-Search (IJCAI 2022) — 78 citations",
]
recent_pubs = [
    "LazyDiT (AAAI 2025) — 49 citations",
    "BitsFusion (NeurIPS 2024) — 43 citations",
    "TextCraftor (CVPR 2024) — 32 citations",
    "SF-V (NeurIPS 2024) — 31 citations",
    "SnapGen (CVPR 2025) — 21 citations",
    "SnapGen-V (CVPR 2025) — 19 citations",
    "SDA (FPL 2024) — 20 citations",
    "Towards Real-Time Segmentation on the Edge (AAAI 2023) — 25 citations",
]
add_simple_list(slide, Inches(1.0), Inches(1.45), Inches(5.55), Inches(4.9), "Highly Cited Contributions", top_pubs)
add_simple_list(slide, Inches(6.8), Inches(1.45), Inches(5.55), Inches(4.9), "Recent Generative and Systems Work", recent_pubs, accent=ACCENT_GREEN)


# Slide 8: patents and service
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_title(slide, "Patents, Service, and Collaboration", "Professional activities beyond publications.")

add_simple_list(
    slide,
    Inches(1.0),
    Inches(1.45),
    Inches(5.8),
    Inches(4.8),
    "Patents and Applications",
    [
        "EfficientFormer Vision Transformer — US Patent 12,236,668 (2025)",
        "Text-to-Image Diffusion Model Rearchitecture — US Patent 12,469,273 (2025)",
        "Vision Transformer for MobileNet Size and Speed — US Patent App. 18/080,993",
        "Step Distillation for Latent Diffusion Models — US Patent App. 18/434,411",
        "Automatic Image Generation Using Latent Structural Diffusion — US Patent App. 18/429,251",
    ],
    accent=ACCENT_ORANGE,
)
add_simple_list(
    slide,
    Inches(7.1),
    Inches(1.45),
    Inches(5.25),
    Inches(2.25),
    "Professional Service",
    [
        "Conference reviewer: NeurIPS, ICCV, ECCV, CVPR, and related venues.",
        "Journal reviewer: IEEE TCAD, IEEE TCAS, and related journals.",
        "Teaching Assistant: Advances in Deep Learning, Northeastern University.",
    ],
    accent=ACCENT_BLUE,
)
add_simple_list(
    slide,
    Inches(7.1),
    Inches(4.0),
    Inches(5.25),
    Inches(2.25),
    "Collaboration Interests",
    [
        "Efficient AI and on-device foundation models",
        "Generative vision systems under resource constraints",
        "Hardware-aware model design and acceleration",
    ],
    accent=ACCENT_GREEN,
)


# Slide 9: thank you
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BLUE)

tb = add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.0))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Thank You"
set_font(r, size=46, bold=True, color=WHITE)
add_paragraph(tf, "Open to research collaborations and discussions on efficient AI and generative modeling.", size=18, color=RGBColor(0xB7, 0xCB, 0xE3), align=PP_ALIGN.CENTER, space_after=Pt(24))
add_paragraph(tf, "yanyu.li@northeastern.edu", size=17, color=RGBColor(0xB7, 0xCB, 0xE3), align=PP_ALIGN.CENTER, space_after=Pt(10))
add_paragraph(tf, "Google Scholar: scholar.google.com/citations?user=XUj8koUAAAAJ", size=17, color=RGBColor(0xB7, 0xCB, 0xE3), align=PP_ALIGN.CENTER, space_after=Pt(10))
add_paragraph(tf, "GitHub: github.com/liyy201912", size=17, color=RGBColor(0xB7, 0xCB, 0xE3), align=PP_ALIGN.CENTER, space_after=Pt(0))


out_path = "/Users/yli16/projects/yy-webpage/presentation.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
